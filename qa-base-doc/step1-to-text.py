import os
import re
import markdown
import docx
import openpyxl
import pptx
from bs4 import BeautifulSoup

################################################################################
### Step 1
################################################################################
    

def parse_file(file_path, save_text_folder, parser):
    text = parser(file_path)
    # 获取文件名称
    file_name = os.path.basename(file_path)
    # 拼接text文件路径
    text_file_path = os.path.join(save_text_folder, f'{file_name}.txt')
    # 打开text文件并写入解析后的文本内容
    with open(text_file_path, 'w', encoding='utf-8') as f:
        f.write(text)

def remove_center_tag(soup):
    # 移除center标签及其内容
    for center_tag in soup.find_all('center'):
        center_tag.extract()

# def parse_html_file(html_file_path, save_text_folder):
#     # 读取html文件内容
#     with open(html_file_path, 'r', encoding='utf-8') as f:
#         html_content = f.read()
#     # 使用BeautifulSoup解析html文件
#     soup = BeautifulSoup(html_content, 'html.parser')
#     # 去除center标签（印象笔记导出后在center标签中包含了一大串不可读的字符）
#     remove_center_tag(soup)
#     # 获取html文件名称
#     html_file_name = os.path.basename(html_file_path)
#     # 去除html后缀
#     html_file_name = re.sub(r'\.html$', '', html_file_name)
#     # 拼接text文件路径
#     text_file_path = os.path.join(save_text_folder, f'{html_file_name}.txt')
#     # 打开text文件并写入解析后的文本内容
#     with open(text_file_path, 'w', encoding='utf-8') as f:
#         f.write(soup.get_text())

def html_parser(file_path, content=''):
    # 读取文件内容
    if len(content) == 0:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    # 使用BeautifulSoup解析html文件
    soup = BeautifulSoup(content, 'html.parser')
    # 去除center标签（印象笔记导出后在center标签中包含了一大串不可读的字符）
    remove_center_tag(soup)
    return soup.get_text()

def markdown_parser(file_path):
    # 读取文件内容
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    text = markdown.markdown(content)
    return html_parser(file_path, text)

def docx_parser(file_path):
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def xlsx_parser(file_path):
    wb = openpyxl.load_workbook(file_path)
    # 获取workbook中所有的sheet名字
    sheet_names = wb.sheetnames
    full_text = []
    # 循环读取每个sheet的内容
    for sheet_name in sheet_names:
        # 根据sheet名字获取sheet对象
        ws = wb[sheet_name]
        for row in ws.values:
            for cell in row:
                if cell is not None and isinstance(cell, str):
                    full_text.append(cell)
    return '\n'.join(full_text)

def pptx_parser(file_path):
    prs = pptx.Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                full_text.append(shape.text)
    return '\n'.join(full_text)


def parse_folder(html_folder_path, save_text_folder):
    # 遍历html文件夹下的所有.html文件
    for root, _, files in os.walk(html_folder_path):
        for file in files:
            # 获取文件路径
            file_path = os.path.join(root, file)
            print(f'提取{file_path}到文本...')
            if file.endswith('.html'):
                # 解析html文件并写入文本文件
                parse_file(file_path, save_text_folder, html_parser)
            elif file.endswith('.md'):
                # 解析md文件并写入文本文件
                parse_file(file_path, save_text_folder, markdown_parser)     
            elif file.endswith('.docx') or file.endswith('.doc'):
                # 解析doc文件并写入文本文件
                parse_file(file_path, save_text_folder, docx_parser)   
            elif file.endswith('.xlsx'):
                # 解析xlsx文件并写入文本文件
                parse_file(file_path, save_text_folder, xlsx_parser)   
            elif file.endswith('.pptx'):
                # 解析pptx文件并写入文本文件
                parse_file(file_path, save_text_folder, pptx_parser)       
                
# 读取doc目录下的文件，提取其中的文本，写入text文件夹下
parse_folder('doc', 'text')
print('---------------------------text生成完毕，写入text文件夹')
